# -*- coding: utf-8 -*-
"""MiniChess AI slide deck with code panels (Traditional Chinese)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY_DK = "0F1B3D"
NAVY    = "1E2761"
BLUE    = "3B5BA9"
ICE     = "CADCFC"
AMBER   = "F2A900"
BG      = "F4F6FB"
WHITE   = "FFFFFF"
TEXT    = "1A2238"
MUTED   = "5B6577"
BORDER  = "DDE3EF"
TINT    = "EEF2FB"
RED_T   = "FBECEC"
RED     = "C0392B"
GRN_T   = "EAF5EE"
GRN     = "2E7D4F"
# code panel
CBG = "12203F"
CFG = "E6ECF7"
CCM = "86C2A6"   # comment green
CHI = "FFC24B"   # highlight amber
CKW = "8FB6FF"   # keyword blue

FONT = "Microsoft JhengHei"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor.from_string(bg)
    return s

def set_ea(run, font_name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", font_name)

def set_mono(run):
    rPr = run._r.get_or_add_rPr()
    for tag, fnt in (("a:latin", MONO), ("a:cs", MONO), ("a:ea", FONT)):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", fnt)

def text(s, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        if "sa" in para: p.space_after = Pt(para["sa"])
        if "line" in para: p.line_spacing = para["line"]
        for seg in para["segs"]:
            txt, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            r = p.add_run(); r.text = txt; r.font.size = Pt(size)
            r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
            set_ea(r, seg[5] if len(seg) > 5 else FONT)
    return tb

def card(s, l, t, w, h, fill=WHITE, border=BORDER, radius=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if border:
        sp.line.color.rgb = RGBColor.from_string(border); sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp

def dot(s, l, t, d, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def badge(s, l, t, d, fill, label, color=WHITE, size=18):
    sp = dot(s, l, t, d, fill)
    tf = sp.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(size)
    r.font.bold = True; r.font.color.rgb = RGBColor.from_string(color)
    set_ea(r, FONT)
    return sp

def title_bar(s, kicker, title):
    text(s, 0.65, 0.42, 12, 0.4, [{"segs": [(kicker, 14, AMBER, True)]}])
    text(s, 0.62, 0.74, 12.1, 0.8, [{"segs": [(title, 28, NAVY, True)]}])

def code_panel(s, l, t, w, h, lines, fs=11.5):
    card(s, l, t, w, h, fill=CBG, border=None, radius=0.03)
    tb = s.shapes.add_textbox(Inches(l + 0.28), Inches(t + 0.2), Inches(w - 0.56), Inches(h - 0.4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12
        if isinstance(line, str): line = [(line, CFG)]
        elif isinstance(line, tuple): line = [line]
        for seg in line:
            txt = seg[0]; col = seg[1] if len(seg) > 1 else CFG
            r = p.add_run(); r.text = txt; r.font.size = Pt(fs)
            r.font.color.rgb = RGBColor.from_string(col); set_mono(r)
    return tb

def expl(s, summary, points, x=7.9, y=1.95, w=4.9):
    paras = [{"segs": [(summary, 15.5, NAVY, True)], "sa": 11, "line": 1.22}]
    for pt in points:
        paras.append({"segs": [("•  ", 13, AMBER, True), (pt, 13.5, TEXT)], "sa": 8, "line": 1.18})
    text(s, x, y, w, 4.7, paras)

def code_slide(kicker, title, lines, summary, points, fs=11.5):
    s = slide()
    title_bar(s, kicker, title)
    code_panel(s, 0.6, 1.78, 7.0, 4.75, lines, fs=fs)
    expl(s, summary, points)
    return s

# ===================================================== S1 Title
s = slide(NAVY_DK)
text(s, 0.9, 1.05, 11, 0.5, [{"segs": [("Mini Project 2 · 迷你西洋棋 AI", 16, ICE, True)]}])
text(s, 0.88, 2.0, 11.6, 1.5, [{"segs": [("MiniChess AI 演算法說明", 46, WHITE, True)]}])
text(s, 0.9, 3.45, 11.6, 0.9, [{"segs": [("從 Negamax 到 PVS 主變例搜尋引擎(含程式碼解說)", 23, ICE)]}])
cx = 0.9
for c in ["vs Weak  100%", "vs Strong  100%", "vs Boss  100%"]:
    card(s, cx, 4.65, 3.25, 0.75, fill=AMBER, border=None, radius=0.5)
    text(s, cx, 4.65, 3.25, 0.75, [{"segs": [(c, 18, NAVY, True)]}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += 3.55
text(s, 0.9, 6.55, 11.5, 0.4,
     [{"segs": [("PVS + 置換表 + 靜止搜尋 + 選擇性深度  |  10 局全勝(2000ms,雙方輪流先手)", 13, ICE)]}])

# ===================================================== S2 Overview
s = slide()
title_bar(s, "OVERVIEW", "整體架構:分層強化")
text(s, 0.65, 1.85, 6.0, 4.6, [
    {"segs": [("引擎以 ", 16, TEXT), ("Negamax + 迭代加深", 16, NAVY, True), (" 為骨架,", 16, TEXT)], "sa": 10, "line": 1.25},
    {"segs": [("再一層層往上疊加:核心剪枝 → 加速結構 → 選擇性深度 → 評估函數。", 16, TEXT)], "sa": 16, "line": 1.25},
    {"segs": [("每一層都讓搜尋「更快」或「更準」,", 15, MUTED)], "sa": 4, "line": 1.25},
    {"segs": [("最終在相同時間內看得更深、下得更強。", 15, MUTED)], "line": 1.25},
])
layers = [
    ("評估函數", "子力 · 位置表 · 王安全 · 兵形", ICE, NAVY),
    ("選擇性深度", "將軍延伸 · LMR · 空著裁剪 · 期望窗口", BLUE, WHITE),
    ("加速結構", "置換表 · 走法排序 · 靜止搜尋", NAVY, WHITE),
    ("核心搜尋", "Negamax → Alpha-Beta → PVS", NAVY_DK, WHITE),
]
ly = 1.8
for name, desc, fill, col in layers:
    card(s, 7.0, ly, 5.7, 1.05, fill=fill, border=None, radius=0.1)
    text(s, 7.3, ly + 0.13, 5.1, 0.45, [{"segs": [(name, 18, col, True)]}])
    text(s, 7.3, ly + 0.56, 5.1, 0.4, [{"segs": [(desc, 12.5, col)]}])
    ly += 1.18

# ===================================================== S3 Negamax
code_slide("CORE SEARCH 1/3", "Negamax 負極大值", [
    [("if", CKW), (" (game_state == WIN)", CFG)],
    [("    return", CKW), (" P_MAX - ply;", CFG), ("        // 越快將死越好", CCM)],
    "",
    [("int", CKW), (" best = M_MAX;", CFG)],
    [("for", CKW), (" (auto& a : legal_actions) {", CFG)],
    [("    State* nx = next_state(a);", CFG)],
    [("    int score = ", CFG), ("-eval_ctx(nx, depth-1, ...)", CHI), (";", CFG)],
    [("    if", CKW), (" (score > best) best = score;", CFG), ("  // 取最大", CCM)],
    "}",
    [("return", CKW), (" best;", CFG)],
],
    "統一雙方視角:取最大,並對子節點取負。",
    ["一條公式處理雙方(零和賽局)",
     "score = −eval_ctx(child) 是核心",
     "正分 = 對走子方有利",
     "終端:可吃王回傳 P_MAX − ply",
     "是 Alpha-Beta / PVS 的基礎"])
# fix first line: combine the leading two segs into one paragraph
# (handled by passing list where needed)

# ===================================================== S4 Alpha-Beta
code_slide("CORE SEARCH 2/3", "Alpha-Beta 剪枝", [
    [("for", CKW), (" (auto& a : moves) {", CFG), ("       // 已 MVV-LVA 排序", CCM)],
    [("  int score =", CFG), (" -eval_ctx(nx, d-1,", CFG)],
    [("                    ", CFG), ("-beta, -alpha", CHI), (", ...);", CFG)],
    "",
    [("  if", CKW), (" (score > best)  best  = score;", CFG)],
    [("  if", CKW), (" (score > alpha) alpha = score;", CFG)],
    [("  if", CKW), (" (alpha >= beta) ", CFG), ("break;", CHI), ("   // beta cutoff", CCM)],
    "}",
],
    "用 [α, β] 窗口剪掉對手不會允許的分枝。",
    ["α:我方已保證的下限",
     "β:對手會容許的上限",
     "score ≥ β → 對手不會走 → 剪枝",
     "結果與完整 minimax 相同",
     "最佳排序下約 O(b^(d/2)),深度加倍",
     "搭配 MVV-LVA 排序提高剪枝率"])

# ===================================================== S5 PVS
code_slide("CORE SEARCH 3/3", "PVS 主變例搜尋", [
    [("if", CKW), (" (first) {", CFG), ("                 // 主變例", CCM)],
    [("  score = -eval(nx, d-1, ", CFG), ("-beta, -alpha", CHI), (");", CFG)],
    [("} ", CFG), ("else", CKW), (" {", CFG), ("                  // 其餘走法", CCM)],
    [("  score = -eval(nx, d-1, ", CFG), ("-alpha-1, -alpha", CHI), (");", CFG)],
    [("  ", CFG), ("// 零寬窗口意外更優 → 完整窗口重搜", CCM)],
    [("  if", CKW), (" (score > alpha && score < beta)", CFG)],
    [("    score = -eval(nx, d-1, ", CFG), ("-beta, -alpha", CHI), (");", CFG)],
    "}",
],
    "第一手用完整窗口,其餘用零寬窗口快篩。",
    ["零寬窗口 [α, α+1]:只問「是否更好」",
     "布林測試比完整搜尋快很多",
     "排序良好時,重搜很少發生",
     "因此比一般 alpha-beta 更快",
     "實際對局採用此演算法"])

# ===================================================== S6 Move ordering
code_slide("MOVE ORDERING", "走法排序:讓剪枝有效", [
    [("// 吃子:MVV-LVA;殺手:-1;安靜:-10", CCM)],
    [("int", CKW), (" as = av", CFG)],
    [("  ? (", CFG), ("piece_val[av]*10 - piece_val[aa]", CHI), (")", CFG)],
    [("  : (a_kill ? ", CFG), ("-1", CHI), (" : ", CFG), ("-10", CHI), (");", CFG)],
    "",
    [("// 根節點先試置換表走法", CCM)],
    [("rotate(moves, tt_move);", CFG)],
],
    "好棋先試,剪枝才會有效。",
    ["順序:置換表手 → 吃子 → 殺手 → 安靜",
     "MVV-LVA:吃大子、用小子吃",
     "殺手:曾造成 cutoff 的安靜走法",
     "置換表手:前一層算出的最佳手",
     "好棋先試 → 更早 cutoff → 看更深"])

# ===================================================== S7 Quiescence
code_slide("QUIESCENCE", "靜止搜尋:避免視界效應", [
    [("int", CKW), (" stand_pat = evaluate(...);", CFG)],
    [("if", CKW), (" (stand_pat >= beta) ", CFG), ("return", CKW), (" beta;", CFG)],
    [("if", CKW), (" (stand_pat > alpha) alpha = stand_pat;", CFG)],
    "",
    [("// delta 剪枝:吃了也拉不回 alpha 就跳過", CCM)],
    [("if", CKW), (" (stand_pat + capture_val[cap] + 50", CFG)],
    [("                       < alpha) ", CFG), ("continue;", CHI)],
    "",
    [("// 只搜吃子,直到局面安靜", CCM)],
],
    "葉節點只追吃子到安靜才評估。",
    ["視界效應:停在被吃子前會誤判佔優",
     "stand-pat:「不吃」的評估基準",
     "delta pruning:跳過無望的吃子",
     "上限 MAX_QDEPTH = 8",
     "讓葉節點的評估在戰術上可靠"])

# ===================================================== S8 TT + Zobrist
code_slide("TRANSPOSITION TABLE", "置換表 與 Zobrist 雜湊", [
    [("// 探測:同局面直接重用", CCM)],
    [("TTEntry& e = g_tt[key & TT_MASK];", CFG)],
    [("if", CKW), (" (e.key == key && e.depth >= depth) {", CFG)],
    [("  if", CKW), (" (e.flag == TT_EXACT) ", CFG), ("return", CKW), (" e.score;", CFG)],
    "}",
    "",
    [("// Zobrist 增量更新(走子時)", CCM)],
    [("h ^= zob[p][orig][from];", CFG), ("   // 移出", CCM)],
    [("h ^= zob[p][moved][to];", CFG), ("    // 移入", CCM)],
],
    "快取局面結果;雜湊 O(1) 增量更新。",
    ["100 萬筆(TT_BITS = 20)雜湊表",
     "key 驗證碰撞;深度優先替換",
     "存深度、分數、界限、最佳走法",
     "Zobrist:每(子,格,色)一隨機數",
     "將死分數依 ply 正規化"])

# ===================================================== S9 Selective depth
code_slide("SELECTIVE DEPTH", "選擇性深度:延伸 / 縮減 / 裁剪", [
    [("// 將軍延伸:強制應對多搜 1 層", CCM)],
    [("int", CKW), (" ext = (ply<20 && nx->is_king_attacked())", CFG)],
    [("          ? ", CFG), ("1", CHI), (" : ", CFG), ("0", CHI), (";", CFG)],
    "",
    [("// 後期走法縮減 (LMR)", CCM)],
    [("if", CKW), (" (!ext && depth>=3 && move_index>=3", CFG)],
    [("           && !capture) reduction = ", CFG), ("1", CHI), (";", CFG)],
    "",
    [("// 空著裁剪", CCM)],
    [("if", CKW), (" (null_score >= beta) ", CFG), ("return", CKW), (" beta;", CFG)],
],
    "把搜尋深度花在關鍵的著法上。",
    ["將軍延伸:被將軍的線多搜 1 層",
     "LMR:後段安靜手先搜淺,意外才重搜",
     "空著裁剪:跳一手仍 ≥ β 就剪枝",
     "期望窗口:窄窗口 ±50 提高剪枝率",
     "皆為啟發式,設條件控制風險"])

# ===================================================== S10 Evaluation
code_slide("EVALUATION", "評估函數:為局面打分", [
    [("self_score += kp_material[piece];", CFG), ("      // 材料", CCM)],
    [("self_score += pst[piece-1][row][col];", CFG), (" // 位置表", CCM)],
    [("self_score += king_tropism(...);", CFG), ("      // 王接近", CCM)],
    [("//  + 通路兵 / 車開放線 / 兵形罰分", CCM)],
    "",
    [("bonus += ", CFG), ("3 * (self_mob - oppn_mob)", CHI), (";", CFG)],
    "",
    [("return", CKW), (" self_score - oppn_score + bonus;", CFG)],
],
    "分數 = 我方總分 − 對方總分(走子方視角)。",
    ["材料:兵20 馬60 象65 車100 后200",
     "位置表 PST:每子每格的位置加分",
     "王接近度:子力靠近敵王給壓力",
     "通路兵 / 車開放線 / 兵形(疊兵孤兵扣分)",
     "機動力:可走步數差 × 3"])

# ===================================================== S11 The fix
s = slide()
title_bar(s, "THE KEY FIX", "關鍵修復:根節點存入置換表")
code_panel(s, 0.6, 1.78, 7.0, 3.55, [
    [("// 每層結束:把根節點最佳手寫入 TT", CCM)],
    [("if", CKW), (" (best_move != Move{} && !ctx.stop) {", CFG)],
    [("  TTEntry& rt = g_tt[root_key & TT_MASK];", CFG)],
    [("  rt.key       = root_key;", CFG)],
    [("  rt.best_move = ", CFG), ("pack_move(best_move)", CHI), (";", CFG)],
    [("  rt.depth     = depth;", CFG)],
    [("  rt.flag      = TT_EXACT;", CFG)],
    "}",
])
expl(s, "限時被中斷時,不再走出敗著。", [
    "舊版:每層用吃子排序,把「會輸的吃子」排第一",
    "被 kill 取最後回報手 → 開局第 7 手被將死",
    "新版:每層先搜、先回報真正最佳手",
    "固定深度本來就對 → 補上「限時」這條路",
], y=1.95)
card(s, 0.6, 5.55, 7.0, 0.95, fill=NAVY, border=None, radius=0.1)
text(s, 0.6, 5.55, 7.0, 0.95,
     [{"segs": [("結果:Weak  ", 16, WHITE), ("40%", 18, "FF8A80", True),
                ("  →  ", 15, ICE), ("100%", 19, AMBER, True)]}],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===================================================== S12 Results
s = slide()
title_bar(s, "RESULTS", "成果:對三個對手全勝")
text(s, 0.65, 1.5, 12, 0.4, [{"segs": [("測試設定:每步 2000ms、各 10 局、雙方輪流先手", 14, MUTED)]}])
res = [("vs Weak", "40%"), ("vs Strong", "90%"), ("vs Boss", "70%")]
xs = [0.6, 4.74, 8.88]
for i, (name, before) in enumerate(res):
    x = xs[i]
    card(s, x, 2.1, 3.84, 3.7, fill=WHITE)
    text(s, x, 2.45, 3.84, 0.5, [{"segs": [(name, 19, NAVY, True)]}], align=PP_ALIGN.CENTER)
    text(s, x, 3.0, 3.84, 1.3, [{"segs": [("100%", 60, AMBER, True)]}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, 4.45, 3.84, 0.5, [{"segs": [("10 勝 0 敗", 15, GRN, True)]}], align=PP_ALIGN.CENTER)
    card(s, x + 0.55, 5.05, 2.74, 0.6, fill=TINT, border=None, radius=0.3)
    text(s, x + 0.55, 5.05, 2.74, 0.6,
         [{"segs": [("修正前 ", 13, MUTED), (before, 14, RED, True), (" → 100%", 13, MUTED)]}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.65, 6.2, 12, 0.5,
     [{"segs": [("關鍵修正(根節點存入 TT)讓三個對手皆達穩定全勝。", 14, MUTED)]}], align=PP_ALIGN.CENTER)

# ===================================================== S13 Summary
s = slide(NAVY_DK)
text(s, 0.9, 0.85, 11, 0.5, [{"segs": [("SUMMARY", 14, AMBER, True)]}])
text(s, 0.88, 1.3, 11.6, 0.9, [{"segs": [("總結", 36, WHITE, True)]}])
card(s, 0.9, 2.5, 11.5, 1.7, fill=NAVY, border=None, radius=0.06)
text(s, 1.25, 2.5, 10.8, 1.7, [
    {"segs": [("這是一個以 Negamax 為基礎、採用迭代加深的 PVS 引擎:", 17, WHITE, True)], "sa": 6, "line": 1.3},
    {"segs": [("alpha-beta + PVS 為核心,搭配置換表與 MVV-LVA/殺手排序讓剪枝有效,", 16, ICE)], "sa": 6, "line": 1.3},
    {"segs": [("靜止搜尋避免戰術失誤,再以延伸/縮減/裁剪把深度花在關鍵處,餵給完整評估函數。", 16, ICE)], "line": 1.3},
])
takeaways = [("PVS 架構", "比 alpha-beta 更快收斂"),
             ("置換表 + 排序", "讓搜尋看得更深"),
             ("根節點存 TT", "修掉限時敗著:40%→100%")]
tx = 0.9
for name, desc in takeaways:
    card(s, tx, 4.55, 3.7, 1.4, fill=NAVY, border=None, radius=0.1)
    text(s, tx + 0.3, 4.78, 3.2, 0.5, [{"segs": [(name, 16, AMBER, True)]}])
    text(s, tx + 0.3, 5.28, 3.2, 0.55, [{"segs": [(desc, 13, WHITE)], "line": 1.15}])
    tx += 3.92
text(s, 0.9, 6.35, 11.5, 0.7, [{"segs": [("謝謝聆聽 · Weak / Strong / Boss 全數 10:0", 16, ICE, True)]}])

prs.save("MiniChess_AI_簡報_含程式碼.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")
